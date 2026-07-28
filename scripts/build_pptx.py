#!/usr/bin/env python3
"""Build the FastERP walkthrough PPTX from its slide-oriented Markdown."""
from __future__ import annotations

import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

AMBER = RGBColor(0xD9, 0x77, 0x06)
DARK = RGBColor(0x27, 0x21, 0x1A)
MUTE = RGBColor(0x7C, 0x71, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def clean(text):
    return re.sub(r"[*`]", "", re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)).strip()


def parse(path):
    source = open(path, encoding="utf-8").read()
    slides = []
    for part in re.split(r"(?m)^---\s*$", source):
        title, image, bullets, paragraphs = "", None, [], []
        for line in part.splitlines():
            text = line.strip()
            if text.startswith(":::"):
                continue
            if text.startswith("#") and not title:
                title = clean(text.lstrip("# "))
            elif text.startswith("!["):
                match = re.search(r"\(([^)]+)\)", text)
                image = match.group(1) if match else None
            elif re.match(r"^[-*]\s", text):
                bullets.append(clean(text[2:]))
            elif re.match(r"^\d+\.\s", text):
                bullets.append(clean(re.sub(r"^\d+\.\s", "", text)))
            elif text and not text.startswith("```"):
                paragraphs.append(clean(text))
        if title:
            slides.append((title, image, bullets, paragraphs))
    return slides


def textbox(slide, left, top, width, height, lines, size=16, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def build(source, output):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    base = os.path.dirname(os.path.abspath(source))
    estonian = source.endswith("_ee.md")
    footer_text = ("FastERP · Sünteetiline raamatupidamise tööruum · 2026-07-28"
                   if estonian else
                   "FastERP · Synthetic accounting workspace · 2026-07-28")
    for index, (title, image, bullets, paragraphs) in enumerate(parse(source)):
        slide = prs.slides.add_slide(blank)
        if index == 0:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = AMBER
            subtitle = ("Raamatupidamise tööruumi kasutusjuhend"
                        if estonian else "Accounting Workspace User Guide")
            tagline = ("Iseseisev sünteetiline näidis"
                       if estonian else "Self-contained synthetic demonstration")
            box = textbox(slide, Inches(.75), Inches(.45), Inches(11.8), Inches(1.45),
                          [title, subtitle, tagline], 18, WHITE)
            for p in box.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            box.text_frame.paragraphs[0].font.size = Pt(34)
            box.text_frame.paragraphs[0].font.bold = True
            if image and os.path.isfile(os.path.join(base, image)):
                slide.shapes.add_picture(os.path.join(base, image), Inches(3.0), Inches(2.0),
                                         width=Inches(7.33))
            continue
        band = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(.92))
        band.fill.solid()
        band.fill.fore_color.rgb = AMBER
        band.line.fill.background()
        title_box = textbox(slide, Inches(.45), Inches(.15), Inches(12.4), Inches(.55),
                            [title], 25, WHITE)
        title_box.text_frame.paragraphs[0].font.bold = True
        body = bullets or paragraphs[:6]
        has_image = image and os.path.isfile(os.path.join(base, image))
        width = Inches(5.3) if has_image else Inches(12)
        textbox(slide, Inches(.55), Inches(1.25), width, Inches(5.5),
                [("• " + x) for x in bullets] if bullets else paragraphs[:6], 15)
        if has_image:
            slide.shapes.add_picture(os.path.join(base, image), Inches(6.15), Inches(1.35),
                                     width=Inches(6.65))
        footer = textbox(slide, Inches(.5), Inches(7.08), Inches(12.3), Inches(.25),
                         [footer_text], 8, MUTE)
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    prs.save(output)
    print(f"Built {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        raise SystemExit("usage: build_pptx.py SOURCE.md OUTPUT.pptx")
    build(sys.argv[1], sys.argv[2])
